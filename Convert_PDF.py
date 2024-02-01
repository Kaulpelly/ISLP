import os
import fitz
from pptx import Presentation
from pptx.util import Inches

pdf_directory = 'D:\Scripts\ISLP\Slide Notes'
pptx_directory = 'D:\Scripts\ISLP\PPTX Files'


# Get all PDF files in the directory
pdf_files = [file for file in os.listdir(pdf_directory) if file.endswith('.pdf')]

for pdf_file in pdf_files:
    # Create a new PPTX presentation
    pptx_file = os.path.join(pptx_directory, os.path.splitext(pdf_file)[0] + '.pptx')
    presentation = Presentation()

    # Open the PDF file
    pdf_path = os.path.join(pdf_directory, pdf_file)
    doc = fitz.open(pdf_path)

    # Convert each page to an image and add it as a slide in the PPTX presentation
    for page in doc:
        image_path = os.path.splitext(pdf_path)[0] + f'_page{page.number}.png'
        pixmap = page.get_pixmap(dpi=300)  # Set the DPI to 300 for higher resolution
        pixmap.save(image_path, "png")  # Save the pixmap as a PNG file
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=presentation.slide_width, height=presentation.slide_height)

    # Save the PPTX presentation
    presentation.save(pptx_file)

    # Close the PDF file
    doc.close()

